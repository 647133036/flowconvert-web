import sys
import json
import os
import re

import requests

USER_AGENT = "Mozilla/5.0 (Windows NT 10.0; Win64; x64) FlowConvert/1.0"

# ISO 639-2/3 (3-letter) -> ISO 639-1 (2-letter) mapping for detected languages
_ISO3_TO_2 = {
    "zho": "zh", "eng": "en", "jpn": "ja", "kor": "ko", "rus": "ru",
    "ara": "ar", "tha": "th", "vie": "vi", "fra": "fr", "deu": "de",
    "spa": "es", "por": "pt", "ita": "it", "nld": "nl", "pol": "pl",
    "tur": "tr",
}


class NoTranslatorError(Exception):
    pass


# ── Language helpers ──

APP_CODES = {
    "zh": "中文", "en": "英语", "ja": "日语", "ko": "韩语",
    "fr": "法语", "de": "德语", "es": "西班牙语", "pt": "葡萄牙语",
    "ru": "俄语", "ar": "阿拉伯语", "th": "泰语", "vi": "越南语",
    "it": "意大利语", "nl": "荷兰语", "pl": "波兰语", "tr": "土耳其语",
}


def detect_language(text):
    """Heuristic language detection, returns app code or None."""
    s = text.strip()[:5000]
    if not s:
        return None
    han = len(re.findall(r"[\u4e00-\u9fff\u3400-\u4dbf]", s))
    hiragana = len(re.findall(r"[\u3040-\u309f]", s))
    katakana = len(re.findall(r"[\u30a0-\u30ff]", s))
    hangul = len(re.findall(r"[\uac00-\ud7af]", s))
    cyrillic = len(re.findall(r"[\u0400-\u04ff]", s))
    thai = len(re.findall(r"[\u0e00-\u0e7f]", s))
    arabic = len(re.findall(r"[\u0600-\u06ff]", s))
    latin = len(re.findall(r"[A-Za-z]", s))
    total = len([c for c in s if not c.isspace()]) or 1

    if (hiragana + katakana) / total > 0.15:
        return "ja"
    if han / total > 0.15:
        return "zh"
    if hangul / total > 0.15:
        return "ko"
    if cyrillic / total > 0.15:
        return "ru"
    if thai / total > 0.15:
        return "th"
    if arabic / total > 0.15:
        return "ar"
    return None


def translate_chunk(text, source, target):
    """Translate one chunk via translatepy; fallback to MyMemory.
    Returns (translated_text, engine, detected_source_or_None)."""
    src = "auto" if source in (None, "", "auto") else source
    detected = None

    # translatepy (aggregates multiple engines including Google)
    try:
        from translatepy import Translator
        t = Translator()
        result = t.translate(text, target, source_language=src)
        out = result.result
        if out and out.strip() and out.strip() != text.strip():
            engine = "translatepy"
            # Extract auto-detected source language when src was "auto"
            if src in (None, "", "auto") and hasattr(result, "source_language"):
                raw = str(result.source_language)
                detected = _ISO3_TO_2.get(raw, raw)
            return out, engine, detected
    except Exception:
        pass

    # MyMemory fallback (only when source is explicit, not auto)
    try:
        if src not in (None, "", "auto"):
            mm_src = src
            if mm_src == "zh":
                mm_src = "zh-CN"
            langpair = "|".join(x for x in [mm_src, target] if x)
            r = requests.post(
                "https://api.mymemory.translated.net/get",
                data={"q": text, "langpair": langpair},
                timeout=10,
                headers={"User-Agent": USER_AGENT},
            )
            d = r.json()
            if d.get("responseStatus") == 200:
                out = d["responseData"]["translatedText"]
                if out and out.strip() and out.strip() != text.strip():
                    return out, "mymemory", None
    except Exception:
        pass

    raise NoTranslatorError("翻译引擎暂时不可用，请稍后重试")


def split_blocks(text, max_chars=1800):
    """Split text into blocks along paragraph boundaries, each <= max_chars."""
    paragraphs = re.split(r"(?<=\n)", text)
    blocks = []
    cur = ""
    for p in paragraphs:
        if len(cur) + len(p) <= max_chars:
            cur += p
        else:
            if cur.strip():
                blocks.append(cur)
            while len(p) > max_chars:
                cut = p[:max_chars]
                blocks.append(cut)
                p = p[max_chars:]
            cur = p
    if cur.strip() or (text.endswith("\n") and cur == "\n"):
        blocks.append(cur)
    return blocks


def translate_text(text, source, target):
    if not text.strip():
        return text, None, None
    detected = None
    if source == "auto" or source in (None, ""):
        detected = detect_language(text)
        # For Latin-script texts (FR/ES/IT/DE etc.), detect_language returns None;
        # pass "auto" to translatepy so it auto-detects via Google/Microsoft.
        src = detected if detected else "auto"
        source = src
    else:
        src = source

    blocks = split_blocks(text, 1200)
    out = []
    used_engine = None
    for b in blocks:
        if not b.strip():
            out.append(b)
            continue
        try:
            translated, engine, chunk_detected = translate_chunk(b, src, target)
            used_engine = engine
            if chunk_detected and not detected:
                detected = chunk_detected
            out.append(translated)
        except NoTranslatorError:
            out.append(b)
    return "".join(out), detected, used_engine


# ── Document translation ──

def translate_txt(src_path, out_path, source, target):
    with open(src_path, "r", encoding="utf-8", errors="replace") as f:
        text = f.read()
    translated, _, _ = translate_text(text, source, target)
    with open(out_path, "w", encoding="utf-8") as f:
        f.write(translated)
    return out_path


def translate_html(src_path, out_path, source, target):
    from bs4 import BeautifulSoup
    with open(src_path, "r", encoding="utf-8", errors="replace") as f:
        text = f.read()
    soup = BeautifulSoup(text, "lxml")
    for elem in soup.find_all(string=True):
        if not elem.strip():
            continue
        parent = elem.parent
        if parent is not None and parent.name in ("script", "style", "code", "pre", "svg"):
            continue
        try:
            translated, _, _ = translate_text(str(elem), source, target)
            if translated and translated != str(elem):
                elem.replace_with(translated)
        except NoTranslatorError:
            continue
    with open(out_path, "w", encoding="utf-8") as f:
        f.write(str(soup))
    return out_path


def translate_docx(src_path, out_path, source, target):
    import docx
    doc = docx.Document(src_path)
    for para in doc.paragraphs:
        if para.text and para.text.strip():
            try:
                translated, _, _ = translate_text(para.text, source, target)
                if translated and translated != para.text:
                    for run in para.runs:
                        run.text = ""
                    if para.runs:
                        para.runs[0].text = translated
                    else:
                        para.add_run(translated)
            except NoTranslatorError:
                continue
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    if para.text and para.text.strip():
                        try:
                            translated, _, _ = translate_text(para.text, source, target)
                            if translated and translated != para.text:
                                for run in para.runs:
                                    run.text = ""
                                if para.runs:
                                    para.runs[0].text = translated
                                else:
                                    para.add_run(translated)
                        except NoTranslatorError:
                            continue
    doc.save(out_path)
    return out_path


def translate_xlsx(src_path, out_path, source, target, max_cells=500):
    from openpyxl import load_workbook
    wb = load_workbook(src_path)
    count = 0
    for ws in wb.worksheets:
        for row in ws.iter_rows():
            for cell in row:
                if cell.value is not None and isinstance(cell.value, str) and cell.value.strip():
                    try:
                        translated, _, _ = translate_text(cell.value, source, target)
                        if translated and translated != cell.value:
                            cell.value = translated
                            count += 1
                            if count >= max_cells:
                                break
                    except NoTranslatorError:
                        pass
                if count >= max_cells:
                    break
            if count >= max_cells:
                break
        if count >= max_cells:
            break
    wb.save(out_path)
    return out_path


def translate_pptx(src_path, out_path, source, target, max_shapes=100):
    from pptx import Presentation
    prs = Presentation(src_path)
    count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = "".join(run.text for run in para.runs)
                if text and text.strip():
                    try:
                        translated, _, _ = translate_text(text, source, target)
                        if translated and translated != text:
                            for run in para.runs:
                                run.text = ""
                            if para.runs:
                                para.runs[0].text = translated
                            count += 1
                            if count >= max_shapes:
                                break
                    except NoTranslatorError:
                        pass
            if count >= max_shapes:
                break
        if count >= max_shapes:
            break
    prs.save(out_path)
    return out_path


def find_cjk_font():
    candidates = [
        "/usr/share/fonts/truetype/wqy/wqy-zenhei.ttc",
        "/usr/share/fonts/truetype/droid/DroidSansFallbackFull.ttf",
        "/usr/share/fonts/truetype/wqy/wqy-microhei.ttc",
    ]
    for p in candidates:
        if os.path.exists(p):
            return p
    try:
        out = os.popen("fc-list :lang=zh file 2>/dev/null | head -1").read().split(":")[0].strip()
        if out and os.path.exists(out):
            return out
    except Exception:
        pass
    return None


def translate_pdf(src_path, out_path, source, target):
    """OCR 识别 PDF → 翻译 → 生成新 PDF（支持图片/PDF 扫描件）"""
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont

    # 注册 CJK 字体
    cjk_font_path = find_cjk_font()
    font_name = "CJKFont"
    if cjk_font_path:
        try:
            pdfmetrics.registerFont(TTFont(font_name, cjk_font_path))
        except Exception:
            font_name = "Helvetica"
    else:
        font_name = "Helvetica"

    # 尝试 OCR 方案
    try:
        from pdf2image import convert_from_path
        import pytesseract

        images = convert_from_path(src_path, dpi=300)
        if not images:
            raise RuntimeError("PDF 无法转换为图像")

        translated_lines = []
        for i, img in enumerate(images):
            custom_config = r'--oem 3 --psm 6 lang=chi_sim+eng'
            text = pytesseract.image_to_string(img, config=custom_config)
            translated, _, _ = translate_text(text, source, target)
            translated_lines.append(translated)

        full_text = "\n\n".join(translated_lines)

    except (ImportError, RuntimeError):
        # 回退：pdfminer 提取文本
        try:
            from pdfminer.high_level import extract_text
            text = extract_text(src_path)
            translated, _, _ = translate_text(text, source, target)
            full_text = translated
        except Exception:
            full_text = ""

    # 生成 PDF
    from reportlab.lib.pagesizes import A4
    from reportlab.pdfgen import canvas as pdfc

    c = pdfc.Canvas(out_path, pagesize=A4)
    page_w, page_h = A4
    margin = 50
    line_h = 14
    y = page_h - margin
    c.setFont(font_name, 11)

    for line in full_text.splitlines():
        if y < margin:
            c.showPage()
            y = page_h - margin
        c.drawString(margin, y, line if line.strip() else " ")
        y -= line_h

    c.save()
    return out_path
    fallback = out_path.rsplit(".", 1)[0] + ".txt"
    with open(fallback, "w", encoding="utf-8") as f:
        f.write("翻译结果（PDF 转译失败，输出为文本）：\n\n")
        f.write(text if 'text' in dir() else "")
    return fallback


TRANSLATORS = {
    "txt": translate_txt,
    "html": translate_html,
    "htm": translate_html,
    "docx": translate_docx,
    "xlsx": translate_xlsx,
    "pptx": translate_pptx,
    "pdf": translate_pdf,
}


def translate_file(src_path, out_path, source, target):
    ext = src_path.rsplit(".", 1)[-1].lower()
    func = TRANSLATORS.get(ext)
    if func is None:
        raise RuntimeError("不支持的文件格式: .%s" % ext)
    return func(src_path, out_path, source, target)


if __name__ == "__main__":
    mode = sys.argv[1]
    if mode == "text":
        source = sys.argv[2]
        target = sys.argv[3]
        with open(sys.argv[4], "r", encoding="utf-8") as f:
            payload = json.load(f)
        text = payload.get("text", "")
        translated, detected, engine = translate_text(text, source, target)
        print(json.dumps({"text": translated, "detected": detected, "engine": engine}))
    elif mode == "file":
        src = sys.argv[2]
        out = sys.argv[3]
        source = sys.argv[4]
        target = sys.argv[5]
        result = translate_file(src, out, source, target)
        print(json.dumps({"output": result}))
    else:
        raise SystemExit("unknown mode")
